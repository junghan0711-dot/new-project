from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


BASE = Path("/Users/junghanchiu/2026 Data base/Templates/Apple風格簡報-範本.pptx")
OUT_DIR = Path("/Users/junghanchiu/Documents/New project/outputs/blueprint_event_ops")
OUTPUT = OUT_DIR / "2026好時光靚市集_分組工作拆解方法教學_Apple風格_v1.pptx"

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(29, 29, 31)
GRAY = RGBColor(134, 134, 139)
LIGHT = RGBColor(245, 245, 247)
BLUE = RGBColor(0, 113, 227)
SOFT_BLUE = RGBColor(232, 243, 255)
FONT = "PingFang TC"


prs = Presentation(BASE)
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)


def clear_slides(presentation):
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        r_id = slide_id.rId
        presentation.part.drop_rel(r_id)
        slide_id_list.remove(slide_id)


clear_slides(prs)
BLANK = prs.slide_layouts[0]


def add_slide(bg=WHITE):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def set_font(run):
    run.font.name = FONT
    r_pr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        font_node = r_pr.find(qn(f"a:{tag}"))
        if font_node is None:
            font_node = OxmlElement(f"a:{tag}")
            r_pr.append(font_node)
        font_node.set("typeface", FONT)


def text_box(slide, text, x, y, w, h, size=28, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def centered(slide, text, y, size=48, color=WHITE, bold=True):
    return text_box(slide, text, 1.2, y, 10.9, 1.5, size=size, color=color, bold=bold, align=PP_ALIGN.CENTER)


def title(slide, text, y=0.82, size=44, dark=False):
    return text_box(slide, text, 0.95, y, 11.5, 1.05, size=size, color=WHITE if dark else TEXT, bold=True)


def subtitle(slide, text, y=1.82, size=21, dark=False):
    return text_box(slide, text, 0.98, y, 11.1, 0.75, size=size, color=RGBColor(190, 190, 195) if dark else GRAY)


def mini_label(slide, text, x, y, dark=False):
    return text_box(slide, text, x, y, 4.5, 0.35, size=15, color=RGBColor(190, 190, 195) if dark else GRAY, bold=True)


def card(slide, x, y, w, h, mark, heading, body, fill=LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    text_box(slide, mark, x + 0.25, y + 0.22, 0.9, 0.45, size=23, color=BLUE, bold=True)
    text_box(slide, heading, x + 0.25, y + 0.76, w - 0.5, 0.45, size=22, color=TEXT, bold=True)
    text_box(slide, body, x + 0.25, y + 1.33, w - 0.5, h - 1.42, size=16, color=GRAY)


def row(slide, y, left, middle, right):
    text_box(slide, left, 0.95, y, 2.25, 0.42, size=18, color=BLUE, bold=True)
    text_box(slide, middle, 3.0, y, 3.75, 0.42, size=18, color=TEXT, bold=True)
    text_box(slide, right, 6.9, y, 5.3, 0.62, size=17, color=GRAY)


def method_header(slide, method, title_text, example):
    mini_label(slide, method, 0.95, 0.78)
    title(slide, title_text, y=1.18, size=42)
    subtitle(slide, example, y=2.1, size=21)


slide = add_slide(BLACK)
centered(slide, "把工作想細，\n不是靠提醒。", 1.85, size=54)
centered(slide, "是靠方法。", 4.1, size=54)
centered(slide, "2026好時光・靚市集｜分組工作拆解教學", 5.75, size=20, color=RGBColor(190, 190, 195), bold=False)

slide = add_slide(BLACK)
centered(slide, "我們今天要練的，\n不是填表。", 1.85, size=52)
centered(slide, "是把腦中的工作，整理成別人可以接手的工作。", 4.15, size=24, color=RGBColor(190, 190, 195), bold=False)

slide = add_slide(WHITE)
title(slide, "五種方法，讓細節浮出來")
card(slide, 0.75, 2.15, 2.35, 3.05, "01", "KJ 法", "先大量寫，再分類命名。")
card(slide, 3.25, 2.15, 2.35, 3.05, "02", "情境走讀", "從使用者的一天走一遍。")
card(slide, 5.75, 2.15, 2.35, 3.05, "03", "倒推法", "從活動日往前推到今天。")
card(slide, 8.25, 2.15, 2.35, 3.05, "04", "風險預演", "先假設出包，再找防線。")
card(slide, 10.75, 2.15, 1.85, 3.05, "05", "交付物與依賴", "確認誰等誰、交什麼。")

slide = add_slide(WHITE)
method_header(slide, "方法 01｜KJ 法", "先把腦中的細節倒出來", "例：活動組不要只寫「舞台組」，先寫出所有可能發生的小事。")
card(slide, 0.95, 3.05, 3.7, 2.55, "寫", "一張便利貼一件事", "主持稿版本、表演團隊進場、音響測試、貴賓稱謂、開幕合照、摸彩 cue 點。", fill=SOFT_BLUE)
card(slide, 4.85, 3.05, 3.7, 2.55, "分", "把相似事情群聚", "舞台流程、貴賓接待、媒體畫面、現場彩排、突發狀況。")
card(slide, 8.75, 3.05, 3.7, 2.55, "命名", "每一群變成工作包", "例如：舞台彩排包、貴賓接待包、主持稿審核包。")

slide = add_slide(WHITE)
title(slide, "KJ 法現場怎麼帶")
subtitle(slide, "15 分鐘就能看出每組漏了什麼。")
row(slide, 2.75, "3 分鐘", "安靜寫", "每人至少寫 8-10 張，不討論、不批評。")
row(slide, 3.35, "5 分鐘", "貼上牆面", "同類貼在一起，把模糊詞拆開。")
row(slide, 3.95, "5 分鐘", "分群命名", "每一群命名成工作包，而不是抽象概念。")
row(slide, 4.55, "2 分鐘", "選前三項", "每組挑出最容易出包、最需要跨組協調的三群。")

slide = add_slide(WHITE)
method_header(slide, "方法 02｜情境走讀", "從人的路線，找出工作細節", "例：招商組用「攤商到場」走一次，就會發現停車、卸貨、報到、撤場都要拆。")
row(slide, 3.0, "攤商 10:00 到", "車停哪裡？", "誰給停車資訊，停車費怎麼處理，卸貨可停多久。")
row(slide, 3.65, "進場報到", "找誰報到？", "資料袋、攤位圖、營收表、設備需求確認。")
row(slide, 4.3, "找不到位置", "誰處理？", "現場聯絡人、備用攤位圖、對講機通報。")
row(slide, 4.95, "撤場", "怎麼離場？", "垃圾、設備、營收表回收、車輛動線。")

slide = add_slide(WHITE)
title(slide, "四種視角都要走讀")
card(slide, 0.95, 2.35, 2.7, 2.75, "民眾", "來逛的人", "入口、導覽、打卡、兌獎、排隊、廁所、離場。")
card(slide, 3.85, 2.35, 2.7, 2.75, "攤商", "來擺攤的人", "通知、進場、卸貨、設備、營收、撤場。")
card(slide, 6.75, 2.35, 2.7, 2.75, "長官", "來視察的人", "抵達、致詞、合照、巡禮、媒體聯訪。")
card(slide, 9.65, 2.35, 2.7, 2.75, "媒體", "來採訪的人", "採訪角度、畫面、受訪者、新聞資料包。")

slide = add_slide(WHITE)
method_header(slide, "方法 03｜倒推法", "從活動當天往前推", "例：服務台如果 10/3 要順利兌獎，就不能 10/2 才想規則。")
row(slide, 2.95, "10/3", "服務台順利兌獎", "人員知道規則，民眾排隊不塞，獎品簽收完整。")
row(slide, 3.55, "9/25", "工作人員演練", "用假情境測試一次：集點、扭蛋、摸彩、分享禮。")
row(slide, 4.15, "9/18", "物料到位", "兌換券、簽收表、指示牌、獎品分類、備用筆。")
row(slide, 4.75, "9/10", "規則定稿", "兌換條件、限制、例外狀況、誰有決定權。")
row(slide, 5.35, "8/30", "流程初版", "畫出民眾從排隊到領獎的動線。")

slide = add_slide(BLACK)
centered(slide, "倒推法的關鍵問題", 1.55, size=48)
centered(slide, "如果那一天要順利，\n今天以前必須完成什麼？", 3.15, size=42)
centered(slide, "不是問什麼時候開始，而是問最晚不能晚於哪一天。", 5.35, size=22, color=RGBColor(190, 190, 195), bold=False)

slide = add_slide(WHITE)
method_header(slide, "方法 04｜風險預演", "先假設出包，再反推防線", "例：服務台爆掉，不一定是人不夠，可能是規則不清、標示不足、獎品保管混亂。")
card(slide, 0.95, 3.0, 3.65, 2.6, "如果", "民眾排隊塞住入口", "可能原因：兌換流程太慢、指示牌不清、服務台位置不對。")
card(slide, 4.85, 3.0, 3.65, 2.6, "所以", "要先設防線", "分流排隊、規則圖卡、服務台前置演練、現場機動人力。")
card(slide, 8.75, 3.0, 3.65, 2.6, "最後", "寫回總控表", "風險、預防措施、發生時誰決策、誰通報。")

slide = add_slide(WHITE)
method_header(slide, "方法 05｜交付物檢查", "沒有交付物，就不算完成", "例：「聯繫表演團隊」不是完成；完成要能交出一包可執行資料。")
row(slide, 2.9, "不是", "聯繫表演團隊", "這只是一個動作，無法判斷是否真的完成。")
row(slide, 3.55, "要變成", "表演團隊確認包", "報價、演出時間、進場時間、設備需求、聯絡人、付款資料。")
row(slide, 4.2, "再補上", "舞台銜接資料", "主持稿 cue 點、彩排時間、走位需求、音控注意事項。")
row(slide, 4.85, "完成標準", "佩妏確認 + 怡文知道", "若要給分署或長官看，也要寫審核節點。")

slide = add_slide(WHITE)
title(slide, "跨組依賴：誰等誰")
subtitle(slide, "每組最後都要列兩欄：我需要什麼，我會交出什麼。")
card(slide, 0.95, 2.6, 3.55, 2.75, "招商組", "給別人", "廠商名單、攤位需求、卸貨需求、商品亮點。")
card(slide, 4.9, 2.6, 3.55, 2.75, "行銷設計", "需要", "廠商故事、主視覺方向、活動亮點、長官審稿時間。")
card(slide, 8.85, 2.6, 3.55, 2.75, "活動組", "接住", "舞台流程、貴賓動線、攝影畫面、現場 cue 點。")

slide = add_slide(BLACK)
centered(slide, "每組今天的任務", 1.7, size=52)
centered(slide, "不是把工作變多，\n是把模糊的工作變清楚。", 3.45, size=38)
centered(slide, "清楚了，才知道誰要幫誰。", 5.45, size=22, color=RGBColor(190, 190, 195), bold=False)

slide = add_slide(WHITE)
title(slide, "今天的分組練習")
subtitle(slide, "每組用 25 分鐘，把一項工作拆成可以回填總控表的工作包。")
row(slide, 2.65, "5 分鐘", "KJ 法發散", "每人寫出自己想到的細節，一張一件事。")
row(slide, 3.25, "5 分鐘", "分群命名", "把便利貼整理成 3-5 個工作包。")
row(slide, 3.85, "5 分鐘", "情境走讀", "選一個角色，把活動前中後走一遍。")
row(slide, 4.45, "5 分鐘", "倒推與風險", "標出最晚日期、審核點與可能出包處。")
row(slide, 5.05, "5 分鐘", "回填總控表", "補上交付物、前置依賴、協作人、需協助事項。")

slide = add_slide(WHITE)
title(slide, "每組最後要交出這一頁")
card(slide, 0.95, 2.2, 3.55, 3.25, "一", "三個新增工作包", "例如：攤商進場包、服務台兌獎包、貴賓媒體包。")
card(slide, 4.9, 2.2, 3.55, 3.25, "二", "每包的完成標準", "交付物、審核人、完成日、需要誰協作。")
card(slide, 8.85, 2.2, 3.55, 3.25, "三", "兩個跨組需求", "我需要別組給我什麼；我會影響誰。")

slide = add_slide(BLACK)
centered(slide, "寫到別人可以照著做。", 2.15, size=54)
centered(slide, "這才叫總控。", 4.25, size=46)
centered(slide, "2026好時光・靚市集｜分組工作拆解練習", 5.85, size=20, color=RGBColor(190, 190, 195), bold=False)


for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.space_after = Pt(0)
                paragraph.space_before = Pt(0)

prs.save(OUTPUT)
print(OUTPUT)
