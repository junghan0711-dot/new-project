from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


BASE = Path("/Users/junghanchiu/2026 Data base/Templates/Apple風格簡報-範本.pptx")
OUT_DIR = Path("/Users/junghanchiu/Documents/New project/outputs/blueprint_event_ops")
OUTPUT = OUT_DIR / "2026好時光靚市集_總控表工作會議_Apple風格_v3.pptx"

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(29, 29, 31)
GRAY = RGBColor(134, 134, 139)
LIGHT = RGBColor(245, 245, 247)
BLUE = RGBColor(0, 113, 227)

FONT = "PingFang TC"


prs = Presentation(BASE)
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)


def clear_slides(presentation):
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        r_id = slide_id.rId
        presentation.part.drop_rel(r_id)
        slide_id_list.remove(slide_id)


clear_slides(prs)
BLANK = prs.slide_layouts[0]


def add_slide(bg=WHITE):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def text_box(slide, text, x, y, w, h, size=28, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        font_node = r_pr.find(qn(f"a:{tag}"))
        if font_node is None:
            font_node = OxmlElement(f"a:{tag}")
            r_pr.append(font_node)
        font_node.set("typeface", FONT)
    return box


def title(slide, text, dark=False, y=0.85, size=44):
    return text_box(slide, text, 0.95, y, 11.4, 1.1, size=size, color=WHITE if dark else TEXT, bold=True)


def subtitle(slide, text, dark=False, y=2.0, size=22):
    return text_box(slide, text, 0.98, y, 10.9, 1.1, size=size, color=GRAY if not dark else RGBColor(190, 190, 195))


def centered(slide, text, y, size=46, color=WHITE, bold=True):
    return text_box(slide, text, 1.25, y, 10.8, 1.5, size=size, color=color, bold=bold, align=PP_ALIGN.CENTER)


def small_label(slide, text, x, y, dark=False):
    return text_box(slide, text, x, y, 3.2, 0.35, size=15, color=RGBColor(190, 190, 195) if dark else GRAY, bold=True)


def card(slide, x, y, w, h, number, heading, body):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = LIGHT
    text_box(slide, number, x + 0.28, y + 0.25, 0.7, 0.5, size=24, color=BLUE, bold=True)
    text_box(slide, heading, x + 0.28, y + 0.78, w - 0.55, 0.55, size=24, color=TEXT, bold=True)
    text_box(slide, body, x + 0.28, y + 1.42, w - 0.55, h - 1.55, size=17, color=GRAY)


def decision_row(slide, y, owner, item, ask):
    text_box(slide, owner, 0.95, y, 1.2, 0.42, size=17, color=BLUE, bold=True)
    text_box(slide, item, 2.05, y, 4.0, 0.42, size=18, color=TEXT, bold=True)
    text_box(slide, ask, 6.35, y, 5.8, 0.55, size=17, color=GRAY)


def group_row(slide, y, group, owner, scope):
    text_box(slide, group, 0.95, y, 3.05, 0.5, size=19, color=TEXT, bold=True)
    text_box(slide, owner, 4.25, y, 2.0, 0.5, size=18, color=BLUE, bold=True)
    text_box(slide, scope, 6.45, y, 5.8, 0.68, size=16, color=GRAY)


slide = add_slide(BLACK)
centered(slide, "把總控表變成\n真正能執行的工作", 2.05, size=52)
centered(slide, "2026好時光・靚市集｜今日工作會議", 4.35, size=22, color=RGBColor(190, 190, 195), bold=False)

slide = add_slide(BLACK)
centered(slide, "今天不是報進度。", 1.85, size=56)
centered(slide, "今天要把每一件事拆到能交辦、能審核、能倒推。", 3.65, size=26, color=RGBColor(190, 190, 195), bold=False)

slide = add_slide(BLACK)
small_label(slide, "目前總控表狀態", 0.95, 0.9, dark=True)
text_box(slide, "31", 0.9, 1.35, 3.4, 1.8, size=136, color=WHITE, bold=True)
text_box(slide, "項任務", 4.0, 2.35, 2.2, 0.55, size=28, color=WHITE, bold=True)
text_box(slide, "已完成 1｜進行中 14｜未開始 16", 0.98, 4.25, 6.2, 0.5, size=24, color=RGBColor(190, 190, 195))
text_box(slide, "真正的風險不是任務多，\n是任務還沒有拆到可執行。", 7.3, 1.75, 4.6, 2.4, size=34, color=WHITE, bold=True)

slide = add_slide(WHITE)
title(slide, "今天只處理三件事")
card(slide, 0.95, 2.25, 3.55, 3.25, "01", "人", "各組負責人、協作人、替補人是否明確。")
card(slide, 4.9, 2.25, 3.55, 3.25, "02", "事", "每項任務的交付物、審核點與下一步是否清楚。")
card(slide, 8.85, 2.25, 3.55, 3.25, "03", "時間", "近截止事項是否有倒推時程與決策期限。")

slide = add_slide(BLACK)
centered(slide, "人可以兼任，\n責任不能混在一起。", 1.9, size=52)
centered(slide, "今天先把分組邊界講清楚。", 4.45, size=24, color=RGBColor(190, 190, 195), bold=False)

slide = add_slide(WHITE)
title(slide, "人少，就分五組")
subtitle(slide, "組別少一點，但工作屬性要乾淨。")
card(slide, 0.8, 2.45, 2.35, 2.75, "01", "總控行政與風險", "決策、進度、錢、人力、保險、保全、結案。")
card(slide, 3.25, 2.45, 2.35, 2.75, "02", "場地攤商與物流", "場地、停車、卸貨、攤位、廠商旅程。")
card(slide, 5.7, 2.45, 2.35, 2.75, "03", "典禮節目與貴賓", "開幕、主持、表演、Run Sheet、接待。")
card(slide, 8.15, 2.45, 2.35, 2.75, "04", "行銷媒體與素材", "視覺、社群、新聞稿、攝影、曝光。")
card(slide, 10.6, 2.45, 1.95, 2.75, "05", "民眾服務與促銷互動", "兌獎、扭蛋、摸彩、諮詢、排隊。")

slide = add_slide(WHITE)
title(slide, "用工作屬性，而不是用習慣分組")
subtitle(slide, "每件事先問：它屬於哪一種工作？")
group_row(slide, 2.55, "總控行政與風險", "怡文 / 玟樺", "總控表、會議追蹤、預算、採購、保險、保全、工讀生、結案。")
group_row(slide, 3.2, "場地攤商與物流", "莞婷", "場地配置、停車卸貨、廠商名單、攤位分配、攤商須知、廠商核銷。")
group_row(slide, 3.85, "典禮節目與貴賓", "佩妏", "啟動儀式、主持人、表演團隊、Run Sheet、貴賓名單、合照聯訪。")
group_row(slide, 4.5, "行銷媒體與素材", "玟樺", "主視覺、海報、社群、新聞稿、戶外廣告、攝影、短影音。")
group_row(slide, 5.15, "民眾服務與促銷互動", "莞婷", "服務台、集點、扭蛋、摸彩、分享禮、獎品保管、排隊動線；莞蓉協作/備援。")

slide = add_slide(WHITE)
title(slide, "今天要討論的是邊界")
subtitle(slide, "不是重新分誰比較重要，而是讓主責清楚、交接清楚。")
card(slide, 0.95, 2.35, 3.55, 3.0, "01", "同類工作放一起", "場地與停車卸貨放一起；舞台與主持表演放一起；宣傳與媒體素材放一起。")
card(slide, 4.9, 2.35, 3.55, 3.0, "02", "每件事只有一個主責", "跨組協作可以很多，但不能三組共同模糊負責。")
card(slide, 8.85, 2.35, 3.55, 3.0, "03", "高風險任務要有第二人", "人可以兼任，但保險、停車、兌獎、新聞稿、貴賓接待都要有備援。")

slide = add_slide(WHITE)
title(slide, "先看到兩個工作量風險")
subtitle(slide, "五組制可以精簡組別，但要避免少數人被壓垮。")
decision_row(slide, 2.75, "莞婷", "場地攤商 + 民眾互動", "回到上一版主責規劃，但要確認柏升、怡君、莞蓉、巧筑或 PT 如何分擔。")
decision_row(slide, 3.55, "玟樺", "行政風險 + 行銷設計 + 新聞稿 + 物料", "今天要確認哪些任務需要佩妏、怡文或外部廠商協作。")
decision_row(slide, 4.35, "莞蓉 / 巧筑", "服務台協作與備援", "不是整組主責，而是協助兌獎規則、獎品保管、排隊動線與現場支援。")
decision_row(slide, 5.15, "全體", "高風險任務第二人", "停車卸貨、廠商名單、啟動儀式、新聞稿、兌獎規則、保險保全都要有備援。")

slide = add_slide(WHITE)
title(slide, "再確認誰負責")
subtitle(slide, "若人沒有定，後面的表都只是漂亮的空格。", y=1.75)
decision_row(slide, 2.75, "怡文", "總控行政與風險", "總控表、預算整合、工作手冊、跨組協調、成果報告。")
decision_row(slide, 3.35, "莞婷", "場地攤商與物流 / 民眾服務與促銷互動", "廠商、攤位、停車卸貨、獎品、服務台規則。")
decision_row(slide, 3.95, "佩妏", "典禮節目與貴賓", "啟動儀式、貴賓名單、主持表演、Run Sheet、攝影協作。")
decision_row(slide, 4.55, "玟樺", "行政風險 / 行銷媒體與素材", "保險保全、工讀生、宣傳素材、新聞稿、物料。")
decision_row(slide, 5.15, "莞蓉 / 巧筑 / 柏升 / 怡君 / 昱碩", "第二人與現場協作", "服務台支援、扭蛋商品、攤位組、長官貴賓接待。")

slide = add_slide(WHITE)
title(slide, "本週必須決定")
subtitle(slide, "7/10 到 7/15 的黃燈事項，今天不能只說「再確認」。", y=1.75)
decision_row(slide, 2.7, "怡文", "組織與分工", "各組組長今天確認職掌是否需調整。")
decision_row(slide, 3.25, "佩妏", "啟動儀式 / 貴賓名單", "何時送科長裁示，裁示前要交什麼方案。")
decision_row(slide, 3.8, "玟樺", "保險與保全", "龍邦保全外是否再比價，誰比、何時回。")
decision_row(slide, 4.35, "莞婷", "停車與卸貨", "城市車旅資訊、廠商繳費協助流程、卸貨動線。")
decision_row(slide, 4.9, "佩妏", "表演與主持", "茶山飛鄒鼓後，中場表演與主持人聯繫期限。")
decision_row(slide, 5.45, "莞婷", "25 家廠商名單", "7/14 篩選標準、7/15 送分署承辦檢閱版本。")

slide = add_slide(BLACK)
centered(slide, "每一個「待確認」\n都要變成一個日期。", 2.0, size=50)
centered(slide, "不然它就會在活動前一週變成風險。", 4.35, size=24, color=RGBColor(190, 190, 195), bold=False)

slide = add_slide(WHITE)
title(slide, "跨組依賴要攤開")
card(slide, 0.95, 2.18, 3.55, 3.35, "A", "廠商名單", "影響攤位配置、宣傳素材、新聞稿故事、停車卸貨。")
card(slide, 4.9, 2.18, 3.55, 3.35, "B", "啟動儀式", "影響舞台流程、貴賓接待、主持稿、媒體畫面。")
card(slide, 8.85, 2.18, 3.55, 3.35, "C", "場地動線", "影響救護站、服務台、攤商進撤場、人潮疏散。")

slide = add_slide(WHITE)
title(slide, "新聞稿只是範例")
subtitle(slide, "重點不是新聞稿本身，而是每項工作都要用同樣方式拆。", y=1.75)
card(slide, 0.95, 2.55, 2.75, 2.65, "1", "日期", "什麼時候初稿、送審、定稿、發布。")
card(slide, 3.95, 2.55, 2.75, 2.65, "2", "角色", "誰主責、誰協作、誰最後審核。")
card(slide, 6.95, 2.55, 2.75, 2.65, "3", "交付", "最後要交什麼檔案、表單或決議。")
card(slide, 9.95, 2.55, 2.45, 2.65, "4", "風險", "如果延誤，會卡住誰。")

slide = add_slide(WHITE)
title(slide, "會後 24 小時內")
subtitle(slide, "每組回到總控表，補齊這七欄。", y=1.75)
text_box(slide, "下一步", 1.0, 2.65, 2.0, 0.5, size=28, color=TEXT, bold=True)
text_box(slide, "交付物", 4.0, 2.65, 2.0, 0.5, size=28, color=TEXT, bold=True)
text_box(slide, "審核節點", 7.0, 2.65, 2.3, 0.5, size=28, color=TEXT, bold=True)
text_box(slide, "前置依賴", 10.0, 2.65, 2.3, 0.5, size=28, color=TEXT, bold=True)
text_box(slide, "需協助｜最晚完成日｜延誤影響", 1.0, 4.05, 10.8, 0.8, size=34, color=BLUE, bold=True)

slide = add_slide(BLACK)
centered(slide, "總控表不是紀錄。", 1.8, size=54)
centered(slide, "它是讓每個人知道下一步怎麼做的工作系統。", 3.55, size=28, color=RGBColor(190, 190, 195), bold=False)
centered(slide, "今天，請把它補到可以執行。", 5.05, size=24, color=WHITE, bold=True)


for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.space_after = Pt(0)
                paragraph.space_before = Pt(0)

prs.save(OUTPUT)
print(OUTPUT)
